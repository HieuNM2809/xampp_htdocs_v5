"""
Sinh file mẫu cho toàn bộ định dạng MarkItDown hỗ trợ (`02_batch_convert.py`).

Phủ 13 đuôi: docx, xlsx, pptx, pdf, xml, txt, md, png, jpg, gif, bmp, wav, zip.

Cách dùng:
    python make_samples.py

Các đuôi không sinh (yêu cầu môi trường ngoài):
    .doc / .ppt : format binary cũ, không có lib write phổ biến.
    .msg        : Outlook proprietary.
    .mp3 / .m4a : cần ffmpeg / encoder.
    .epub       : zip + XML schema phức tạp, ít test value.
"""
import sys
import wave
import zipfile
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")

SAMPLES = Path(__file__).parent / "samples"
SAMPLES.mkdir(exist_ok=True)


def make_xml():
    (SAMPLES / "sample.xml").write_text(
        """<?xml version="1.0" encoding="UTF-8"?>
<catalog>
  <product id="P001">
    <name>Sữa rửa mặt Cetaphil 250ml</name>
    <category>Skincare</category>
    <price currency="VND">295000</price>
  </product>
  <product id="P002">
    <name>Kem chống nắng Anessa 60ml</name>
    <category>Sunscreen</category>
    <price currency="VND">650000</price>
  </product>
</catalog>
""",
        encoding="utf-8",
    )


def make_txt():
    (SAMPLES / "sample.txt").write_text(
        """Báo cáo nhanh Q1/2026
=====================

- Doanh thu online vượt cửa hàng lần đầu tiên.
- Mobile App tăng trưởng 78%.
- Skincare dẫn đầu danh mục.

Liên hệ: bi@hasaki.vn
""",
        encoding="utf-8",
    )


def make_md():
    (SAMPLES / "sample.md").write_text(
        """# Ghi chú tuần 22

## Việc đã làm
- Setup MarkItDown demo
- Test batch convert

## TODO
- [ ] Đóng gói thành package
- [x] Viết README
""",
        encoding="utf-8",
    )


def make_docx():
    from docx import Document
    doc = Document()
    doc.add_heading("Báo cáo doanh thu Q1/2026", 0)
    doc.add_paragraph(
        "Tổng quan kết quả kinh doanh quý 1 năm 2026 của Hasaki. "
        "Số liệu chưa kiểm toán, chỉ dùng nội bộ."
    )
    doc.add_heading("Doanh thu theo kênh", 1)
    table = doc.add_table(rows=4, cols=3)
    table.style = "Light Grid"
    headers = ["Kênh", "Doanh thu (triệu VND)", "Tăng trưởng"]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
    for i, row in enumerate(
        [("Cửa hàng", "1.250", "+12%"), ("Online", "2.340", "+45%"), ("App", "890", "+78%")],
        start=1,
    ):
        for j, val in enumerate(row):
            table.cell(i, j).text = val
    doc.add_paragraph("Mảng Skincare vẫn dẫn đầu, chiếm 42% tổng doanh thu.")
    doc.save(SAMPLES / "sample.docx")


def make_xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Doanh thu"
    ws1.append(["Kênh", "Doanh thu", "Tăng trưởng"])
    for row in [("Cửa hàng", 1250, "+12%"), ("Online", 2340, "+45%"), ("App", 890, "+78%")]:
        ws1.append(row)
    ws2 = wb.create_sheet("Sản phẩm")
    ws2.append(["Mã", "Tên", "Giá VND"])
    ws2.append(["P001", "Cetaphil 250ml", 295000])
    ws2.append(["P002", "Anessa 60ml", 650000])
    ws2.append(["P003", "The Ordinary Niacinamide", 360000])
    wb.save(SAMPLES / "sample.xlsx")


def make_pptx():
    from pptx import Presentation
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Báo cáo Q1/2026"
    s.placeholders[1].text = "Hasaki Technology — Phòng BI"

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Doanh thu theo kênh"
    tf = s.placeholders[1].text_frame
    tf.text = "Cửa hàng: 1.250 triệu (+12%)"
    tf.add_paragraph().text = "Online: 2.340 triệu (+45%)"
    tf.add_paragraph().text = "App: 890 triệu (+78%)"

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Kết luận"
    s.placeholders[1].text_frame.text = (
        "Online lần đầu vượt cửa hàng. Tiếp tục đầu tư kênh số."
    )
    prs.save(SAMPLES / "sample.pptx")


def make_pdf():
    # reportlab Helvetica không có glyph tiếng Việt → dùng tiếng Anh để PDF render đúng.
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(str(SAMPLES / "sample.pdf"), pagesize=A4)
    c.setFont("Helvetica-Bold", 18)
    c.drawString(72, 780, "Q1/2026 Revenue Report")
    c.setFont("Helvetica", 11)
    y = 740
    for line in [
        "Channel breakdown:",
        "  - Store: 1,250M VND  (+12%)",
        "  - Online: 2,340M VND (+45%)",
        "  - Mobile App: 890M VND (+78%)",
        "",
        "Highlights:",
        "  Online surpassed store for the first time.",
        "  Mobile App grew fastest at +78%.",
    ]:
        c.drawString(72, y, line)
        y -= 18
    c.save()


def make_image(ext: str):
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (480, 120), color=(245, 245, 250))
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("arial.ttf", 22)
    except OSError:
        font = ImageFont.load_default()
    draw.text((24, 38), f"MarkItDown sample {ext}", fill=(20, 20, 60), font=font)
    draw.text((24, 70), "Hasaki BI - Q1/2026", fill=(80, 80, 120), font=font)
    img.save(SAMPLES / f"sample{ext}")


def make_wav():
    # 1 giây silence, 8 kHz mono — đủ để MarkItDown nhận diện format.
    sr = 8000
    with wave.open(str(SAMPLES / "sample.wav"), "wb") as wav:
        wav.setnchannels(1)
        wav.setsampwidth(2)
        wav.setframerate(sr)
        wav.writeframes(b"\x00\x00" * sr)


def make_zip():
    with zipfile.ZipFile(SAMPLES / "sample.zip", "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("readme.txt", "Đây là file trong zip.\nDùng test convert .zip qua MarkItDown.\n")
        z.writestr("data.csv", "name,value\nFoo,1\nBar,2\n")


def main():
    print(f"Generating samples in {SAMPLES}")
    builders = [
        ("xml",  make_xml),
        ("txt",  make_txt),
        ("md",   make_md),
        ("docx", make_docx),
        ("xlsx", make_xlsx),
        ("pptx", make_pptx),
        ("pdf",  make_pdf),
        ("png",  lambda: make_image(".png")),
        ("jpg",  lambda: make_image(".jpg")),
        ("gif",  lambda: make_image(".gif")),
        ("bmp",  lambda: make_image(".bmp")),
        ("wav",  make_wav),
        ("zip",  make_zip),
    ]
    ok = fail = 0
    for ext, fn in builders:
        try:
            fn()
            print(f"  [OK]   sample.{ext}")
            ok += 1
        except Exception as e:
            print(f"  [FAIL] sample.{ext}: {e}")
            fail += 1
    print(f"\n{ok} OK, {fail} FAIL")


if __name__ == "__main__":
    main()
